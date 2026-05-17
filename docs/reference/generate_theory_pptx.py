from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
OUTPUT_PPTX = BASE_DIR / "Theory_Assignments_PPT.pptx"


COLORS = {
    "ink": RGBColor(27, 31, 42),
    "charcoal": RGBColor(38, 45, 57),
    "paper": RGBColor(248, 250, 252),
    "muted": RGBColor(93, 107, 125),
    "line": RGBColor(215, 222, 232),
    "teal": RGBColor(12, 146, 145),
    "amber": RGBColor(241, 157, 56),
    "coral": RGBColor(217, 88, 73),
    "green": RGBColor(69, 154, 109),
    "blue": RGBColor(67, 103, 180),
    "violet": RGBColor(123, 88, 171),
    "white": RGBColor(255, 255, 255),
}


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.blank = self.prs.slide_layouts[6]
        self.slide_no = 0

    def slide(self, title: str | None = None, section: str | None = None, dark: bool = False):
        slide = self.prs.slides.add_slide(self.blank)
        self.slide_no += 1

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLORS["ink"] if dark else COLORS["paper"]

        if title:
            title_color = COLORS["white"] if dark else COLORS["ink"]
            self.text(slide, title, 0.55, 0.28, 8.8, 0.45, size=24, bold=True, color=title_color)
        if section:
            sec_color = COLORS["amber"] if dark else COLORS["teal"]
            self.text(slide, section.upper(), 10.35, 0.32, 2.4, 0.28, size=8.5, bold=True, color=sec_color, align=PP_ALIGN.RIGHT)

        self.footer(slide, dark=dark)
        return slide

    def footer(self, slide, dark: bool = False):
        color = RGBColor(180, 189, 202) if dark else COLORS["muted"]
        self.text(slide, "Devulapalli Tharun | 252CS009 | Theory Assignments", 0.55, 7.04, 6.8, 0.18, size=7.5, color=color)
        self.text(slide, f"{self.slide_no:02d}", 12.25, 7.03, 0.45, 0.18, size=8, bold=True, color=color, align=PP_ALIGN.RIGHT)

    def text(
        self,
        slide,
        text: str,
        x,
        y,
        w,
        h,
        size=14,
        bold=False,
        color=None,
        align=PP_ALIGN.LEFT,
        font_name="Aptos",
        valign=None,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.01)
        tf.margin_bottom = Inches(0.01)
        if valign:
            tf.vertical_anchor = valign

        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or COLORS["ink"]
        return box

    def card(self, slide, x, y, w, h, fill="white", line="line", rounded=True):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = COLORS[fill] if isinstance(fill, str) else fill
        shp.line.color.rgb = COLORS[line] if isinstance(line, str) else line
        shp.line.width = Pt(1.1)
        return shp

    def stat(self, slide, label: str, value: str, x, y, w=2.45, h=0.9, accent="teal"):
        self.card(slide, x, y, w, h)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS[accent]
        bar.line.fill.background()
        self.text(slide, value, x + 0.22, y + 0.12, w - 0.35, 0.25, size=20, bold=True)
        self.text(slide, label, x + 0.22, y + 0.51, w - 0.35, 0.22, size=8.8, color=COLORS["muted"])

    def bullets(self, slide, items: list[str], x, y, w, h, size=12.5, color=None, gap=5):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"- {item}"
            p.level = 0
            p.space_after = Pt(gap)
            p.font.name = "Aptos"
            p.font.size = Pt(size)
            p.font.color.rgb = color or COLORS["ink"]
        return box

    def code(self, slide, code: str, x, y, w, h, size=9.5):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(30, 34, 44)
        shp.line.color.rgb = RGBColor(70, 78, 94)

        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.16)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.12)
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_ANCHOR.TOP

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = code
        run.font.name = "Consolas"
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(235, 239, 245)
        return shp

    def table(self, slide, headers: list[str], rows: list[list[str]], x, y, w, h, font_size=9):
        tbl = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
        for c, header in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["charcoal"]
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = Pt(font_size)
                p.font.color.rgb = COLORS["white"]

        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["white"] if r % 2 else RGBColor(242, 246, 250)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(font_size)
                    p.font.color.rgb = COLORS["ink"]
        return tbl


def add_flow_box(deck: Deck, slide, label: str, x, y, w, h, accent="teal"):
    deck.card(slide, x, y, w, h, fill="white", line=accent)
    deck.text(slide, label, x + 0.14, y + 0.18, w - 0.28, h - 0.36, size=11.5, bold=True, color=COLORS[accent], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_arrow(slide, x, y, w, h, color: RGBColor):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow


def build_theory_deck() -> None:
    d = Deck()

    # 1. Cover with introduction.
    s = d.slide(dark=True)
    d.text(s, "NIT Karnataka, Surathkal", 0.72, 0.42, 5.1, 0.28, size=12, color=COLORS["amber"], bold=True)
    d.text(s, "THEORY SUBMISSION", 0.72, 0.95, 6.8, 0.52, size=30, color=COLORS["white"], bold=True)
    d.text(s, "Assignment 1 and Assignment 2", 0.72, 1.55, 7.2, 0.42, size=22, color=COLORS["white"], bold=True)
    d.text(s, "Natural Language Processing for Data Systems and Management", 0.74, 2.05, 9.3, 0.34, size=15, color=RGBColor(210, 220, 235))

    intro_items = [
        ("Name", "Devulapalli Tharun"),
        ("Roll No", "252CS009"),
        ("Department", "Computer Science (MTech)"),
        ("Institute", "National Institute of Technology Karnataka, Surathkal"),
        ("Academic Year", "2025-2026"),
    ]
    y = 2.85
    for key, value in intro_items:
        d.text(s, key, 0.82, y, 1.9, 0.24, size=9.5, color=RGBColor(170, 184, 203), bold=True)
        d.text(s, value, 2.85, y, 7.7, 0.27, size=10.8, color=COLORS["white"])
        y += 0.47

    d.card(s, 10.55, 1.02, 2.1, 5.5, fill="charcoal", line="muted")
    d.text(s, "Deck Focus", 10.76, 1.26, 1.7, 0.24, size=11, bold=True, color=COLORS["amber"], align=PP_ALIGN.CENTER)
    d.bullets(
        s,
        [
            "Chapter 1 summary",
            "Reference algorithms review",
            "Research gap",
            "Proposed TFTH algorithm",
            "Evaluation and conclusion",
        ],
        10.77,
        1.75,
        1.7,
        3.8,
        size=8.8,
        color=COLORS["white"],
        gap=7,
    )

    # 2. Source and scope.
    s = d.slide("Source Material and Submission Scope", "Overview")
    d.text(s, "This presentation is created directly from your submitted Theory assignment files in this folder.", 0.68, 0.85, 11.2, 0.32, size=13, color=COLORS["muted"])
    d.stat(s, "Assignments covered", "2", 0.72, 1.32, accent="teal")
    d.stat(s, "Target slide count", "15", 3.35, 1.32, accent="amber")
    d.stat(s, "Reference algorithms", "7+", 5.98, 1.32, accent="blue")
    d.stat(s, "Proposed algorithm", "TFTH", 8.61, 1.32, accent="green")
    d.bullets(
        s,
        [
            r"Theory 1\Chapter_1.pdf and Chapter_1.docx: Chapter 1 book draft on NLP in data systems.",
            r"Theory 2\chapter1_algorithms.pdf and chapter1_algorithms.docx: algorithm survey and proposed method.",
            "Assignment 1 coverage in this deck is concise: key chapter ideas and summary points only.",
            "Assignment 2 coverage is deeper, with dedicated slides emphasizing the proposed TF-Transformer Hybrid algorithm.",
        ],
        0.76,
        2.55,
        11.85,
        2.9,
        size=12.2,
        gap=8,
    )
    d.table(
        s,
        ["Deliverable", "What this PPT includes"],
        [
            ["Assignment 1 (Chapter 1)", "Core concepts, chapter structure, enterprise relevance, chapter takeaway"],
            ["Assignment 2 (Algorithms)", "Existing algorithms, limitations, TFTH design, scoring equation, expected gains"],
        ],
        0.78,
        5.7,
        11.8,
        1.2,
        font_size=9.5,
    )

    # 3. Assignment 1 summary map.
    s = d.slide("Assignment 1: Chapter 1 Overview", "Theory 1")
    d.bullets(
        s,
        [
            "Chapter title: Introduction to NLP in Data Systems.",
            "Purpose: explain how NLP converts unstructured text into structured, queryable knowledge.",
            "Core storyline: data growth challenge -> NLP evolution -> data taxonomy -> system integration -> enterprise applications.",
            "Style: combines conceptual discussion with complexity-oriented mathematical framing.",
        ],
        0.78,
        1.0,
        6.15,
        2.6,
        size=13,
        gap=9,
    )
    d.table(
        s,
        ["Chapter 1 section", "Main idea"],
        [
            ["1.1 Data-intensive computing", "Big-data growth and compute/memory constraints"],
            ["1.2 Evolution of NLP", "Rule-based -> statistical -> neural paradigms"],
            ["1.3 Data taxonomy", "Structured vs semi-structured vs unstructured data"],
            ["1.4 NLP in data systems", "Semantic lifting architecture and throughput constraints"],
            ["1.5 Enterprise applications", "Sentiment, IE, retrieval, chatbots, fraud detection"],
            ["1.6 Chapter summary", "Key equations, complexity landscape, next chapter roadmap"],
        ],
        6.75,
        1.06,
        5.78,
        4.35,
        font_size=9.2,
    )
    d.stat(s, "Primary theme", "NLP as semantic bridge", 0.95, 5.32, 3.2, accent="teal")
    d.stat(s, "Data view", "High entropy -> structured", 4.4, 5.32, 3.4, accent="violet")
    d.stat(s, "System focus", "Scalable enterprise deployment", 8.1, 5.32, 3.9, accent="green")

    # 4. Chapter 1 key technical points.
    s = d.slide("Chapter 1: Key Technical Highlights", "Theory 1")
    d.table(
        s,
        ["Theme", "Key point from chapter", "Why it matters in systems"],
        [
            ["Data growth", "Data volumes grow exponentially over time", "Requires scalable storage and pipeline design"],
            ["NLP evolution", "Symbolic -> statistical -> transformer models", "Model choice impacts quality and compute cost"],
            ["Semantic gap", "Unstructured text needs computational lifting to structure", "Enables analytics and SQL-like querying"],
            ["Throughput limits", "Latency and bandwidth constraints bound pipeline speed", "Forces online/offline architecture split"],
            ["Complexity focus", "Attention and retrieval complexity are central bottlenecks", "Guides algorithm and infra decisions"],
        ],
        0.75,
        1.12,
        11.9,
        3.4,
        font_size=9.0,
    )
    d.bullets(
        s,
        [
            "Assignment 1 is strong on conceptual rigor: it frames NLP as part of data engineering, not only language modeling.",
            "The chapter repeatedly links algorithmic complexity to practical deployment constraints in enterprise systems.",
            "This foundation motivates Assignment 2, where algorithm selection and hybrid retrieval design become central.",
        ],
        0.85,
        4.85,
        11.6,
        1.65,
        size=12.2,
        gap=8,
    )

    # 5. Assignment 1 applications summary.
    s = d.slide("Assignment 1: Enterprise NLP Summary", "Theory 1")
    d.table(
        s,
        ["Application", "NLP task", "Operational value"],
        [
            ["Customer analytics", "Sentiment and intent classification", "Faster product and support decisions"],
            ["Knowledge extraction", "NER + relation extraction", "Builds searchable enterprise knowledge graphs"],
            ["Intelligent search", "Semantic retrieval", "Higher relevance than keyword-only search"],
            ["Conversational systems", "Dialogue understanding", "Automates repetitive support workflows"],
            ["Fraud/risk analysis", "Anomaly detection in textual logs", "Improves early warning capabilities"],
        ],
        0.8,
        1.2,
        11.8,
        3.5,
        font_size=9.2,
    )
    d.card(s, 0.85, 5.05, 11.7, 1.55)
    d.text(
        s,
        "Assignment 1 takeaway: NLP should be engineered as a production data-system layer balancing semantic quality, latency, and scalability.",
        1.05,
        5.44,
        11.25,
        0.7,
        size=14,
        bold=True,
        color=COLORS["ink"],
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    # 6. Assignment 2 section header.
    s = d.slide(dark=True)
    d.text(s, "ASSIGNMENT 2", 0.78, 1.05, 5.4, 0.55, size=34, bold=True, color=COLORS["amber"])
    d.text(s, "Study of Existing NLP Algorithms and New Proposal", 0.8, 1.72, 9.8, 0.45, size=22, bold=True, color=COLORS["white"])
    d.text(s, "Focus: Reference algorithm analysis + proposed TF-Transformer Hybrid (TFTH)", 0.82, 2.3, 10.6, 0.42, size=14, color=RGBColor(210, 220, 235))
    d.stat(s, "Algorithm families", "Lexical + neural + ANN", 0.95, 4.25, 3.5, accent="teal")
    d.stat(s, "Proposed method", "TFTH", 4.75, 4.25, 2.6, accent="amber")
    d.stat(s, "Target use case", "Large-scale retrieval", 7.62, 4.25, 3.3, accent="green")

    # 7. Existing algorithm landscape.
    s = d.slide("Assignment 2: Existing Algorithm Landscape", "Theory 2")
    d.table(
        s,
        ["Algorithm", "Pipeline stage", "Strength", "Limitation"],
        [
            ["TF-IDF", "Lexical retrieval", "Fast, interpretable, sparse indexing", "Weak semantic understanding"],
            ["Word2Vec (Skip-Gram)", "Word representation", "Dense semantic vectors", "Context-independent embeddings"],
            ["LSTM", "Sequence modeling", "Handles temporal dependencies", "Sequential training bottleneck"],
            ["Transformer", "Deep contextual encoding", "Parallel and expressive", "Quadratic attention cost"],
            ["HNSW", "ANN retrieval index", "Low-latency nearest neighbor search", "Quality depends on embedding quality"],
            ["BERT", "Pre-trained contextual model", "Strong transfer learning", "Heavy inference cost"],
            ["SBERT", "Sentence embedding retrieval", "Fast similarity search after indexing", "Still compute-heavy at scale"],
        ],
        0.75,
        1.2,
        11.88,
        4.55,
        font_size=8.8,
    )
    d.bullets(
        s,
        [
            "The assignment clearly traces a path from lexical methods to transformer-based semantic retrieval.",
            "No single reference algorithm gives the best trade-off across relevance, speed, and scalability.",
        ],
        0.82,
        5.98,
        11.7,
        0.8,
        size=12,
        gap=6,
    )

    # 8. Classical-to-neural progression.
    s = d.slide("Algorithm Deep Dive: TF-IDF, Word2Vec, LSTM", "Theory 2")
    d.bullets(
        s,
        [
            "TF-IDF gives a robust lexical baseline with cosine similarity over sparse vectors.",
            "Word2Vec learns distributed semantics and captures local context relationships.",
            "LSTM improves sequence handling but introduces sequential processing overhead.",
            "These methods establish critical foundations but are not sufficient alone for modern retrieval quality + speed demands.",
        ],
        0.8,
        1.0,
        6.2,
        2.8,
        size=12.8,
        gap=9,
    )
    d.code(
        s,
        "TF-IDF score(d, q) = cos(x_d, x_q)\n"
        "Word2Vec objective: maximize context prediction\n"
        "LSTM recurrence: h_t = f(h_(t-1), x_t)\n\n"
        "Insight: lexical precision + semantic context + sequence memory\n"
        "still need efficient large-scale retrieval integration.",
        0.88,
        3.95,
        5.95,
        2.15,
        size=9.4,
    )
    d.table(
        s,
        ["Method", "Typical complexity (high level)", "Best use"],
        [
            ["TF-IDF retrieval", "Near-linear sparse ops + top-k ranking", "Keyword-heavy corpora"],
            ["Word2Vec training", "Corpus-scale iterative embedding updates", "Word-level semantic enrichment"],
            ["LSTM inference", "Sequential over token length", "Temporal sequence tasks"],
        ],
        6.95,
        1.15,
        5.55,
        3.05,
        font_size=8.8,
    )
    d.stat(s, "Key message", "Strong foundation, limited alone", 7.05, 4.55, 5.3, accent="violet")

    # 9. Transformer-era algorithms.
    s = d.slide("Algorithm Deep Dive: Transformer, BERT, SBERT, HNSW", "Theory 2")
    d.table(
        s,
        ["Component", "Role in retrieval pipeline", "Trade-off"],
        [
            ["Transformer encoder", "Contextual representation learning", "High quality, high compute"],
            ["BERT", "Bidirectional pre-training + fine-tuning", "Excellent accuracy, slow pairwise scoring"],
            ["SBERT", "Sentence-level embedding for similarity", "Fast search with ANN, embedding cost remains"],
            ["HNSW", "Approximate nearest-neighbor index", "Very low-latency top-k, approximate results"],
        ],
        0.78,
        1.14,
        11.82,
        2.9,
        font_size=9.0,
    )
    d.bullets(
        s,
        [
            "BERT delivers rich semantics but scales poorly for direct pairwise comparisons at large N.",
            "SBERT + HNSW improves retrieval speed significantly by embedding once and indexing.",
            "However, pure semantic retrieval may miss exact lexical intent in specialized data-system queries.",
        ],
        0.84,
        4.35,
        11.65,
        1.45,
        size=12.5,
        gap=8,
    )
    d.card(s, 0.9, 5.95, 11.6, 0.88, fill="white", line="teal")
    d.text(s, "This limitation motivates a hybrid lexical + semantic retrieval strategy.", 1.12, 6.2, 11.2, 0.3, size=14, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)

    # 10. Research gap.
    s = d.slide("Research Gap and Design Goal", "Proposed Method")
    d.table(
        s,
        ["Observed gap", "Impact", "Design target in TFTH"],
        [
            ["Keyword-only retrieval misses semantic context", "Low recall for paraphrased queries", "Add transformer semantic candidates"],
            ["Semantic-only retrieval may ignore exact domain terms", "Precision drop for exact terminology", "Retain TF-IDF lexical candidate path"],
            ["Large-scale corpora need low-latency search", "Slow ranking blocks production usage", "Use HNSW ANN for fast semantic retrieval"],
            ["Single-score methods are rigid", "Poor adaptability across query types", "Weighted fusion score with tunable alpha"],
        ],
        0.78,
        1.14,
        11.85,
        3.15,
        font_size=8.9,
    )
    d.code(
        s,
        "Hybrid principle:\n"
        "C = C_TF-IDF union C_Transformer\n"
        "score(d) = alpha * s_TF-IDF(q, d) + (1 - alpha) * s_TR(q, d)\n\n"
        "Recommended from submission: alpha around 0.3 for balanced behavior.",
        0.88,
        4.55,
        11.55,
        1.85,
        size=10.3,
    )

    # 11. Proposed algorithm architecture.
    s = d.slide("Proposed Algorithm: TF-Transformer Hybrid (TFTH)", "Proposed Method")
    d.text(s, "Architecture combines lexical precision and semantic retrieval in one ranking pipeline.", 0.78, 0.9, 11.8, 0.3, size=13, color=COLORS["muted"])

    add_flow_box(d, s, "Query q", 0.75, 2.05, 1.65, 0.95, accent="teal")
    add_flow_box(d, s, "TF-IDF\nTop-k1 candidates", 2.95, 1.45, 2.4, 1.1, accent="amber")
    add_flow_box(d, s, "Transformer encode +\nHNSW Top-k2", 2.95, 3.15, 2.4, 1.1, accent="blue")
    add_flow_box(d, s, "Candidate union C", 6.05, 2.3, 2.1, 0.95, accent="violet")
    add_flow_box(d, s, "Fusion score\n(alpha blend)", 8.75, 2.3, 2.2, 0.95, accent="green")
    add_flow_box(d, s, "Final ranked\nTop-k results", 11.45, 2.3, 1.45, 0.95, accent="teal")

    add_arrow(s, 2.45, 2.33, 0.45, 0.33, COLORS["muted"])
    add_arrow(s, 5.45, 1.8, 0.45, 0.25, COLORS["muted"])
    add_arrow(s, 5.45, 3.5, 0.45, 0.25, COLORS["muted"])
    add_arrow(s, 8.25, 2.58, 0.45, 0.25, COLORS["muted"])
    add_arrow(s, 10.98, 2.58, 0.4, 0.25, COLORS["muted"])

    d.bullets(
        s,
        [
            "Dual candidate generation increases recall while preserving exact-match precision.",
            "Fusion stage allows query-adaptive balancing through alpha.",
            "HNSW keeps semantic search latency practical for large corpora.",
        ],
        0.85,
        5.25,
        12.0,
        1.35,
        size=12.1,
        gap=8,
    )

    # 12. Proposed pseudocode.
    s = d.slide("TFTH Retrieval Procedure", "Proposed Method")
    d.code(
        s,
        "Algorithm TFTH_Retrieve(Corpus D, Query q, k1, k2, alpha):\n"
        "1. Build TF-IDF index for D\n"
        "2. Encode D with Transformer -> embeddings E\n"
        "3. Build HNSW index over E\n"
        "4. C_tf = TopK_TFIDF(q, k1)\n"
        "5. e_q = TransformerEncode(q)\n"
        "6. C_tr = TopK_HNSW(e_q, k2)\n"
        "7. C = Union(C_tf, C_tr)\n"
        "8. For each d in C:\n"
        "      s(d) = alpha*s_tf(q,d) + (1-alpha)*s_tr(q,d)\n"
        "9. Return TopK documents by s(d)",
        0.85,
        1.2,
        7.2,
        5.2,
        size=10.2,
    )
    d.table(
        s,
        ["Parameter", "Meaning", "Typical choice"],
        [
            ["k1", "Lexical candidate size", "20-100 depending on corpus"],
            ["k2", "Semantic candidate size", "20-100 with ANN budget"],
            ["alpha", "Lexical-vs-semantic weight", "0.3 (recommended start)"],
            ["Index", "Semantic ANN structure", "HNSW"],
        ],
        8.25,
        1.3,
        4.35,
        2.35,
        font_size=9.0,
    )
    d.bullets(
        s,
        [
            "This flow maps directly to your Assignment 2 proposed Algorithm 8 (TFTH Retrieval).",
            "The method is implementation-friendly for enterprise document search systems.",
        ],
        8.3,
        4.15,
        4.2,
        2.0,
        size=11.5,
        gap=9,
    )

    # 13. Expected comparison.
    s = d.slide("Expected Performance Positioning of TFTH", "Proposed Method")
    d.table(
        s,
        ["Method", "Precision (exact terms)", "Semantic recall", "Latency", "Overall retrieval quality"],
        [
            ["TF-IDF only", "High", "Low", "Very fast", "Medium"],
            ["SBERT + HNSW", "Medium", "High", "Fast", "High"],
            ["TFTH (proposed)", "High", "High", "Fast", "Very high (balanced)"],
        ],
        0.8,
        1.12,
        11.85,
        2.25,
        font_size=9.3,
    )
    d.code(
        s,
        "Why TFTH should work better:\n"
        "1. Preserves lexical intent for domain-specific terms.\n"
        "2. Recovers semantically relevant paraphrases.\n"
        "3. Keeps retrieval scalable using ANN indexing.\n"
        "4. Offers tunable behavior through alpha for different workloads.",
        0.9,
        3.75,
        6.0,
        2.3,
        size=10.5,
    )
    d.card(s, 7.2, 3.78, 5.3, 2.25, fill="white", line="green")
    d.text(
        s,
        "Proposed contribution highlighted:\nTFTH is a practical hybrid retrieval algorithm designed for real-world NLP data systems where both relevance and speed matter.",
        7.45,
        4.07,
        4.8,
        1.6,
        size=12.5,
        bold=True,
        color=COLORS["green"],
        align=PP_ALIGN.LEFT,
    )

    # 14. Evaluation plan.
    s = d.slide("Evaluation Plan for Proposed TFTH", "Proposed Method")
    d.table(
        s,
        ["Evaluation aspect", "Plan"],
        [
            ["Datasets", "Use chapter-domain corpora and enterprise-style document sets"],
            ["Baselines", "TF-IDF baseline and SBERT+HNSW semantic baseline"],
            ["Metrics", "Precision@k, Recall@k, nDCG, query latency, index memory"],
            ["Ablation", "Vary alpha, k1, k2 to test lexical/semantic balance sensitivity"],
            ["Outcome goal", "Show better balanced relevance-speed trade-off than single-path methods"],
        ],
        0.8,
        1.15,
        11.85,
        3.1,
        font_size=9.1,
    )
    d.bullets(
        s,
        [
            "Success criterion: TFTH should improve recall over TF-IDF while preserving precision better than semantic-only retrieval.",
            "System criterion: latency should remain production-friendly using ANN-based semantic retrieval.",
            "This positions TFTH as an engineering-oriented contribution, not only a theoretical one.",
        ],
        0.85,
        4.55,
        11.6,
        1.85,
        size=12.3,
        gap=8,
    )

    # 15. Conclusion.
    s = d.slide("Conclusion and Key References", "Closing")
    d.bullets(
        s,
        [
            "Assignment 1 established a strong conceptual foundation for NLP in modern data systems.",
            "Assignment 2 reviewed major reference algorithms from lexical, neural, and ANN retrieval pipelines.",
            "Your proposed TFTH algorithm is the central contribution: a balanced hybrid retrieval strategy.",
            "The final direction is clear: evaluate TFTH against standard baselines and tune alpha for domain behavior.",
        ],
        0.82,
        1.0,
        11.7,
        2.2,
        size=12.8,
        gap=9,
    )
    d.table(
        s,
        ["Reference", "Why used"],
        [
            ["Lafferty et al. (CRF, 2001)", "Sequence labeling foundations"],
            ["Pennington et al. (GloVe, 2014)", "Word representation context"],
            ["Vaswani et al. (Transformer, 2017)", "Core modern sequence architecture"],
            ["Reimers and Gurevych (SBERT, 2019)", "Efficient sentence retrieval baseline"],
        ],
        0.82,
        3.65,
        11.75,
        2.1,
        font_size=9.3,
    )
    d.card(s, 0.85, 5.95, 11.7, 0.86, fill="white", line="teal")
    d.text(s, "Thank you - Theory PPT completed in lab-style standard with 15 slides.", 1.1, 6.2, 11.2, 0.3, size=14, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)

    d.prs.save(OUTPUT_PPTX)
    print(f"Saved polished theory presentation: {OUTPUT_PPTX}")
    print(f"Total slides generated: {len(d.prs.slides)}")


if __name__ == "__main__":
    build_theory_deck()
