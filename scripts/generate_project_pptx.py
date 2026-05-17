from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PPTX = ROOT / "Kalakar_2_Project_Presentation.pptx"


COLORS = {
    "ink": RGBColor(27, 31, 42),
    "charcoal": RGBColor(38, 45, 57),
    "paper": RGBColor(248, 250, 252),
    "muted": RGBColor(93, 107, 125),
    "line": RGBColor(215, 222, 232),
    "teal": RGBColor(12, 146, 145),
    "amber": RGBColor(241, 157, 56),
    "coral": RGBColor(217, 88, 73),
    "green": RGBColor(69, 154, 109),
    "blue": RGBColor(67, 103, 180),
    "violet": RGBColor(123, 88, 171),
    "white": RGBColor(255, 255, 255),
}


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.blank = self.prs.slide_layouts[6]
        self.slide_no = 0

    def slide(self, title: str | None = None, section: str | None = None, dark: bool = False):
        slide = self.prs.slides.add_slide(self.blank)
        self.slide_no += 1
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLORS["ink"] if dark else COLORS["paper"]

        if title:
            self.text(
                slide,
                title,
                0.55,
                0.28,
                9.3,
                0.48,
                size=24,
                bold=True,
                color=COLORS["white"] if dark else COLORS["ink"],
            )
        if section:
            self.text(
                slide,
                section.upper(),
                10.35,
                0.32,
                2.4,
                0.28,
                size=8.5,
                bold=True,
                color=COLORS["amber"] if dark else COLORS["teal"],
                align=PP_ALIGN.RIGHT,
            )
        self.footer(slide, dark=dark)
        return slide

    def footer(self, slide, dark: bool = False):
        color = RGBColor(180, 189, 202) if dark else COLORS["muted"]
        self.text(
            slide,
            "Devulapalli Tharun | 252CS009 | Food Price Forecasting",
            0.55,
            7.04,
            7.0,
            0.18,
            size=7.5,
            color=color,
        )
        self.text(slide, f"{self.slide_no:02d}", 12.25, 7.03, 0.45, 0.18, size=8, bold=True, color=color, align=PP_ALIGN.RIGHT)

    def text(
        self,
        slide,
        text: str,
        x,
        y,
        w,
        h,
        size=14,
        bold=False,
        color=None,
        align=PP_ALIGN.LEFT,
        font_name="Aptos",
        valign=None,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.01)
        tf.margin_bottom = Inches(0.01)
        if valign:
            tf.vertical_anchor = valign
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or COLORS["ink"]
        return box

    def card(self, slide, x, y, w, h, fill="white", line="line", rounded=True):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = COLORS[fill] if isinstance(fill, str) else fill
        shp.line.color.rgb = COLORS[line] if isinstance(line, str) else line
        shp.line.width = Pt(1.05)
        return shp

    def stat(self, slide, label: str, value: str, x, y, w=2.35, h=0.9, accent="teal"):
        self.card(slide, x, y, w, h)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS[accent]
        bar.line.fill.background()
        self.text(slide, value, x + 0.22, y + 0.11, w - 0.35, 0.27, size=19, bold=True)
        self.text(slide, label, x + 0.22, y + 0.51, w - 0.35, 0.24, size=8.7, color=COLORS["muted"])

    def bullets(self, slide, items: list[str], x, y, w, h, size=12.0, color=None, gap=5):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"- {item}"
            p.level = 0
            p.space_after = Pt(gap)
            p.font.name = "Aptos"
            p.font.size = Pt(size)
            p.font.color.rgb = color or COLORS["ink"]
        return box

    def code(self, slide, code: str, x, y, w, h, size=9.2):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(30, 34, 44)
        shp.line.color.rgb = RGBColor(70, 78, 94)
        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.16)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.12)
        tf.margin_bottom = Inches(0.08)
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = code
        run.font.name = "Consolas"
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(235, 239, 245)
        return shp

    def table(self, slide, headers: list[str], rows: list[list[str]], x, y, w, h, font_size=8.8):
        tbl = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
        for c, header in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["charcoal"]
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = Pt(font_size)
                p.font.color.rgb = COLORS["white"]
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["white"] if r % 2 else RGBColor(242, 246, 250)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(font_size)
                    p.font.color.rgb = COLORS["ink"]
        return tbl

    def image(self, slide, path: Path, x, y, w, h, pad=0.06, framed=True):
        if framed:
            self.card(slide, x, y, w, h)
        with Image.open(path) as img:
            iw, ih = img.size
        max_w = w - 2 * pad
        max_h = h - 2 * pad
        scale = min(max_w / iw, max_h / ih)
        draw_w = iw * scale
        draw_h = ih * scale
        px = x + (w - draw_w) / 2
        py = y + (h - draw_h) / 2
        slide.shapes.add_picture(str(path), Inches(px), Inches(py), width=Inches(draw_w), height=Inches(draw_h))


def flow_box(d: Deck, slide, label: str, x, y, w, h, accent="teal", size=10.3):
    d.card(slide, x, y, w, h, fill="white", line=accent)
    d.text(slide, label, x + 0.12, y + 0.12, w - 0.24, h - 0.24, size=size, bold=True, color=COLORS[accent], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def arrow(slide, x, y, w, h, color="muted"):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = COLORS[color]
    shp.line.fill.background()
    return shp


def down_arrow(slide, x, y, w, h, color="muted"):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = COLORS[color]
    shp.line.fill.background()
    return shp


def add_cover(d: Deck):
    s = d.slide(dark=True)
    d.text(s, "NIT Karnataka, Surathkal", 0.72, 0.42, 5.2, 0.28, size=12, color=COLORS["amber"], bold=True)
    d.text(s, "INTERPRETABLE AND UNCERTAINTY-AWARE", 0.72, 0.92, 8.8, 0.45, size=24, color=COLORS["white"], bold=True)
    d.text(s, "FOOD PRICE FORECASTING", 0.72, 1.35, 8.8, 0.52, size=31, color=COLORS["white"], bold=True)
    d.text(s, "for Indian Markets using Temporal Fusion Transformers", 0.74, 1.96, 8.9, 0.36, size=16, color=RGBColor(210, 220, 235))

    details = [
        ("Submitted by", "Devulapalli Tharun"),
        ("Roll No", "252CS009"),
        ("Programme", "CSE - M.Tech"),
        ("Institute", "National Institute of Technology Karnataka, Surathkal"),
        ("Email", "tharun.252cs009@nitk.edu.in"),
        ("Guided by", "Dr. Shashidhar G Koolagudi"),
        ("Internal guide", "L. Preethi"),
    ]
    y = 2.7
    for key, value in details:
        d.text(s, key, 0.82, y, 1.9, 0.23, size=9.3, color=RGBColor(170, 184, 203), bold=True)
        d.text(s, value, 2.85, y, 7.4, 0.26, size=10.6, color=COLORS["white"])
        y += 0.4

    d.card(s, 10.45, 1.0, 2.2, 5.55, fill="charcoal", line="muted")
    emblem = ROOT / "Report&Paper" / "figures" / "NITK_Emblem.png"
    if emblem.exists():
        d.image(s, emblem, 10.84, 1.27, 1.42, 1.42, framed=False)
    d.text(s, "Deck Focus", 10.76, 2.98, 1.68, 0.25, size=11, bold=True, color=COLORS["amber"], align=PP_ALIGN.CENTER)
    d.bullets(
        s,
        [
            "Problem and dataset",
            "Feature pipeline",
            "TFT-XGB fusion",
            "CQR calibration",
            "Results and dashboard",
        ],
        10.75,
        3.45,
        1.75,
        1.72,
        size=8.4,
        color=COLORS["white"],
        gap=3,
    )
    d.text(s, "Best model: TFT-XGBFusion-CQR", 10.72, 5.65, 1.82, 0.45, size=10.2, bold=True, color=COLORS["green"], align=PP_ALIGN.CENTER)


def build_deck() -> None:
    d = Deck()
    add_cover(d)

    # 2
    s = d.slide("Project Question and Design Answer", "Overview")
    d.text(s, "Can we forecast Indian food prices with accuracy, uncertainty, and explanation?", 0.72, 1.02, 10.8, 0.34, size=18, bold=True, color=COLORS["teal"])
    d.bullets(
        s,
        [
            "Point forecast alone is risky during onion and tomato shock periods.",
            "The project predicts q10, q50, and q90 prices for 1-6 months ahead.",
            "TFT gives dynamic feature weights and attention; CQR makes bands empirically reliable.",
            "The final dashboard exposes history, future bands, explanations, and test validation.",
        ],
        0.82,
        1.62,
        7.2,
        2.2,
        size=12.6,
    )
    d.stat(s, "Test MAPE", "11.4%", 8.55, 1.42, accent="teal")
    d.stat(s, "Band coverage", "84.0%", 10.95, 1.42, accent="green")
    d.stat(s, "Test rows", "792", 8.55, 2.55, accent="blue")
    d.stat(s, "Series", "123", 10.95, 2.55, accent="violet")
    d.code(
        s,
        "System answer:\n"
        "WFP prices + NASA weather\n"
        "-> 26 engineered features\n"
        "-> XGBoost point signal\n"
        "-> TFT quantile forecast\n"
        "-> per-commodity CQR\n"
        "-> Streamlit decision dashboard",
        0.82,
        4.05,
        5.25,
        1.85,
        size=10.2,
    )
    d.image(s, ROOT / "Report&Paper" / "figures" / "dashboard_view1_forecast.png", 6.45, 3.82, 6.15, 2.25)

    # 3
    s = d.slide("Dataset and Modeling Scope", "Data")
    d.stat(s, "Raw WFP rows", "145,124", 0.74, 1.05, accent="blue")
    d.stat(s, "Final rows", "16,939", 3.18, 1.05, accent="teal")
    d.stat(s, "Markets", "53", 5.62, 1.05, accent="violet")
    d.stat(s, "Commodities", "3", 8.06, 1.05, accent="amber")
    d.stat(s, "Test period", "2023+", 10.5, 1.05, accent="green")
    d.table(
        s,
        ["Source", "Role", "Selected content"],
        [
            ["WFP food prices", "Target history", "Monthly retail prices, 1994-2026"],
            ["WFP markets", "Spatial join", "53 city coordinates"],
            ["NASA POWER", "Weather covariates", "Rainfall, temperature, humidity"],
            ["Processed panel", "Model input", "Market x commodity monthly series"],
        ],
        0.74,
        2.35,
        6.2,
        2.0,
        font_size=8.7,
    )
    d.table(
        s,
        ["Commodity", "Rows", "Series", "Volatility meaning"],
        [
            ["Onions", "5,003", "38", "Seasonal crisis-prone"],
            ["Tomatoes", "5,245", "45", "Fast spike-prone"],
            ["Rice", "6,691", "40", "Stable staple baseline"],
        ],
        7.25,
        2.35,
        5.35,
        1.55,
        font_size=8.7,
    )
    d.bullets(
        s,
        [
            "Train: up to Dec 2021 for the final TFT family.",
            "Validation: 2022, used for early stopping and CQR offsets.",
            "Test: 2023 onward, held out for final evaluation.",
        ],
        7.38,
        4.25,
        5.0,
        1.2,
        size=11.2,
    )
    d.code(s, "Group key = market + commodity\nTarget = log1p(price Rs/KG)\nHorizon H = 6 months\nEncoder window = 24 months", 0.74, 4.72, 6.2, 1.35, size=10.2)

    # 4
    s = d.slide("End-to-End System Flow", "Pipeline")
    flow_box(d, s, "Raw data\nWFP + NASA", 0.74, 1.15, 1.65, 0.86, "blue")
    arrow(s, 2.48, 1.38, 0.48, 0.34)
    flow_box(d, s, "Cleaning\nand filtering", 3.04, 1.15, 1.62, 0.86, "teal")
    arrow(s, 4.76, 1.38, 0.48, 0.34)
    flow_box(d, s, "Feature panel\n26 columns", 5.32, 1.15, 1.62, 0.86, "violet")
    down_arrow(s, 5.93, 2.1, 0.38, 0.52)
    flow_box(d, s, "XGBoost\npoint baseline", 3.68, 2.78, 1.88, 0.82, "amber")
    flow_box(d, s, "TFT quantiles\nq10 / q50 / q90", 6.1, 2.78, 2.02, 0.82, "teal")
    arrow(s, 5.66, 3.03, 0.36, 0.27, "amber")
    d.text(s, "xgb_log_pred", 5.5, 3.42, 0.82, 0.2, size=7.8, bold=True, color=COLORS["amber"], align=PP_ALIGN.CENTER)
    down_arrow(s, 6.88, 3.72, 0.4, 0.5)
    flow_box(d, s, "Per-commodity\nCQR calibration", 5.62, 4.34, 2.3, 0.86, "green")
    down_arrow(s, 6.57, 5.3, 0.4, 0.52)
    flow_box(d, s, "Streamlit dashboard\nforecast + explain", 5.25, 5.92, 3.05, 0.78, "blue")
    d.card(s, 9.05, 1.18, 3.45, 5.4)
    d.text(s, "Main artifacts", 9.28, 1.46, 2.9, 0.26, size=12.5, bold=True, color=COLORS["teal"])
    d.bullets(
        s,
        [
            "models/xgb_clean_2019.pkl",
            "models/tft_best_xgbfused.ckpt",
            "models/conformal_offsets_step5.json",
            "data/processed/tft_predictions_calibrated_step5.csv",
            "app.py dashboard selects best available family",
        ],
        9.3,
        1.95,
        2.85,
        2.25,
        size=9.4,
    )
    d.code(s, "Final family = step5\nName = TFT-XGBFusion-CQR\nActive when fused checkpoint\nand fused dataset exist.", 9.28, 4.6, 2.9, 1.18, size=8.5)

    # 5
    s = d.slide("Feature Engineering", "Preprocessing")
    d.table(
        s,
        ["Group", "Examples", "Why it matters"],
        [
            ["Static", "commodity, market, state", "Series identity and regional context"],
            ["Calendar", "month, season, sin/cos", "Known future seasonal structure"],
            ["Price memory", "lag_1m, lag_12m, rolling_3m/6m", "Momentum and annual baseline"],
            ["Weather", "rainfall, temperature, humidity", "Supply-side stress signal"],
            ["Shock flags", "rain_deficit, heat_stress", "Nonlinear crisis regimes"],
            ["Fusion", "xgb_log_pred", "Strong point prior for TFT"],
        ],
        0.72,
        1.12,
        6.25,
        3.05,
        font_size=8.5,
    )
    d.code(
        s,
        "log_price = log(1 + price)\n"
        "month_sin = sin(2*pi*month/12)\n"
        "month_cos = cos(2*pi*month/12)\n"
        "roll_3(s,t) = mean(p[s,t], p[s,t-1], p[s,t-2])\n"
        "yoy_change = (p[s,t] - p[s,t-12]) / p[s,t-12]\n"
        "rain_deficit = 1 if rainfall < 25th percentile",
        7.28,
        1.18,
        5.2,
        2.0,
        size=8.9,
    )
    d.bullets(
        s,
        [
            "Features are shared by XGBoost and TFT to keep the comparison controlled.",
            "Log-price reduces spike dominance while preserving multiplicative changes.",
            "Binary weather shocks give the VSN readable, domain-aligned gates.",
            "The fused XGBoost feature is created using a leakage-controlled pre-2020 model.",
        ],
        7.38,
        3.55,
        4.88,
        1.75,
        size=11.1,
    )
    d.text(s, "Simple idea: convert raw price history into a supervised panel where every row knows its season, market, weather state, and price memory.", 0.86, 5.18, 6.0, 0.8, size=14.2, bold=True, color=COLORS["teal"])

    # 6
    s = d.slide("Model Stack: Point + Distribution + Calibration", "Method")
    flow_box(d, s, "XGBoost\n500 trees\nmax_depth=6\neta=0.05", 0.82, 1.28, 2.2, 1.28, "amber", size=9.8)
    flow_box(d, s, "Temporal Fusion\nTransformer\nq10 q50 q90", 3.8, 1.28, 2.28, 1.28, "teal", size=9.8)
    flow_box(d, s, "CQR\nper commodity\nlog-space offset", 6.86, 1.28, 2.28, 1.28, "green", size=9.8)
    flow_box(d, s, "Dashboard\nforecast band\nexplanation", 9.95, 1.28, 2.28, 1.28, "blue", size=9.8)
    arrow(s, 3.14, 1.75, 0.55, 0.34, "muted")
    arrow(s, 6.22, 1.75, 0.55, 0.34, "muted")
    arrow(s, 9.33, 1.75, 0.55, 0.34, "muted")
    d.table(
        s,
        ["Layer", "Output", "Limitation solved"],
        [
            ["XGBoost", "single log-price", "Strong tabular accuracy"],
            ["TFT", "three quantiles per horizon", "Uncertainty + interpretability"],
            ["CQR", "adjusted q10/q90", "Empirical coverage"],
            ["Streamlit", "interactive evidence", "Usability and audit"],
        ],
        0.82,
        3.25,
        5.45,
        2.1,
        font_size=8.7,
    )
    d.code(
        s,
        "Prediction object per future month:\n"
        "lower = q10_calibrated\n"
        "median = q50\n"
        "upper = q90_calibrated\n\n"
        "Risk label uses relative band width:\n"
        "(upper - lower) / median",
        6.78,
        3.22,
        5.35,
        1.7,
        size=9.2,
    )
    d.text(s, "Key design choice: XGBoost is not just averaged with TFT. Its prediction becomes a known covariate, so TFT can learn when to trust or override it.", 6.85, 5.28, 5.2, 0.6, size=12.0, bold=True, color=COLORS["teal"])

    # 7
    s = d.slide("TFT Architecture in This Project", "Method")
    d.card(s, 0.72, 1.05, 11.9, 5.55)
    flow_box(d, s, "Static IDs\ncommodity\nmarket\nstate", 1.0, 1.55, 1.55, 0.95, "blue", size=8.7)
    flow_box(d, s, "Known future\nmonth, season\nxgb_log_pred", 1.0, 3.05, 1.55, 0.95, "green", size=8.7)
    flow_box(d, s, "Observed past\nprice, weather\nlags, shocks", 1.0, 4.55, 1.55, 0.95, "amber", size=8.7)
    flow_box(d, s, "VSN gates\nSoftmax weights", 3.16, 2.12, 1.82, 0.9, "violet", size=9.0)
    flow_box(d, s, "LSTM encoder\n24 months", 5.56, 1.7, 1.72, 0.84, "teal", size=9.0)
    flow_box(d, s, "LSTM decoder\n6 months", 5.56, 3.02, 1.72, 0.84, "teal", size=9.0)
    flow_box(d, s, "Static\nenrichment", 7.86, 2.36, 1.62, 0.86, "blue", size=9.0)
    flow_box(d, s, "Interpretable\nattention", 10.0, 2.36, 1.62, 0.86, "coral", size=9.0)
    flow_box(d, s, "Quantile heads\nq10 q50 q90", 10.0, 4.2, 1.62, 0.86, "green", size=9.0)
    arrow(s, 2.68, 3.18, 0.38, 0.26)
    arrow(s, 5.06, 2.55, 0.4, 0.26)
    arrow(s, 7.36, 2.68, 0.4, 0.26)
    arrow(s, 9.56, 2.68, 0.36, 0.26)
    down_arrow(s, 10.6, 3.36, 0.34, 0.45)
    d.code(
        s,
        "Configuration:\nhidden_size=32\nattention_heads=2\nlstm_layers=1\nhidden_continuous=16\ndropout=0.2\nquantiles=[0.1,0.5,0.9]",
        3.1,
        4.38,
        4.35,
        1.25,
        size=8.6,
    )
    d.bullets(
        s,
        [
            "VSN answers: which input variable matters now?",
            "Attention answers: which past month influenced this horizon?",
            "Quantile heads answer: what is the likely range, not only the center?",
        ],
        7.85,
        5.35,
        3.85,
        0.82,
        size=9.2,
    )

    # 8
    s = d.slide("XGBoost Fusion Without Leakage", "Novelty")
    d.text(s, "The strongest point model is injected into TFT as xgb_log_pred.", 0.76, 1.02, 8.9, 0.34, size=17, bold=True, color=COLORS["teal"])
    d.code(
        s,
        "1. Train auxiliary XGBoost only up to 2019-12.\n"
        "2. Generate xgb_log_pred for all rows after that.\n"
        "3. Add xgb_log_pred as a known real covariate.\n"
        "4. Train TFT up to 2021-12, validate on 2022.\n"
        "5. Test on 2023+ with CQR-calibrated bands.",
        0.82,
        1.72,
        5.45,
        1.75,
        size=9.4,
    )
    d.table(
        s,
        ["Model family", "Training cutoff", "What changed"],
        [
            ["TFT-Base", "2020", "Plain TFT, weaker coverage"],
            ["TFT-Retrain21", "2021", "Fairer recent training window"],
            ["TFT-XGBFusion", "2021", "Adds xgb_log_pred known feature"],
            ["TFT-XGBFusion-CQR", "2021 + 2022 cal", "Final calibrated band"],
        ],
        6.72,
        1.55,
        5.55,
        2.05,
        font_size=8.6,
    )
    flow_box(d, s, "XGB prediction\nstrong local prior", 1.18, 4.45, 2.1, 0.78, "amber")
    arrow(s, 3.42, 4.67, 0.55, 0.3)
    flow_box(d, s, "TFT VSN\nlearns gate weight", 4.1, 4.45, 2.1, 0.78, "violet")
    arrow(s, 6.34, 4.67, 0.55, 0.3)
    flow_box(d, s, "Quantile output\nkeeps uncertainty", 7.02, 4.45, 2.1, 0.78, "teal")
    arrow(s, 9.26, 4.67, 0.55, 0.3)
    flow_box(d, s, "CQR band\nreliability", 9.95, 4.45, 2.1, 0.78, "green")
    d.text(s, "Why pre-2020 XGB? It prevents calibration/test leakage: the signal used inside TFT has not already seen the answer for 2022-2023.", 1.0, 5.78, 10.8, 0.42, size=12.2, bold=True, color=COLORS["coral"])

    # 9
    s = d.slide("Quantile Loss and CQR Calibration", "Math")
    d.code(
        s,
        "Pinball loss for quantile q:\n"
        "QL(y, yhat, q) = q * max(y - yhat, 0)\n"
        "              + (1-q) * max(yhat - y, 0)\n\n"
        "TFT total loss:\n"
        "L = sum over q in {0.1,0.5,0.9}, horizon H=6",
        0.75,
        1.12,
        5.75,
        2.15,
        size=9.2,
    )
    d.code(
        s,
        "CQR score on validation:\n"
        "s_i = max(q10_i - y_i, y_i - q90_i)\n\n"
        "Offset c = quantile_{1-alpha}(s_i)\n"
        "q10_cal = q10 - c\n"
        "q90_cal = q90 + c\n\n"
        "Applied separately for onions, tomatoes, rice.",
        6.85,
        1.12,
        5.45,
        2.15,
        size=9.2,
    )
    d.table(
        s,
        ["Commodity", "CQR offset in log space", "Test coverage"],
        [
            ["Onions", "-0.0277", "90.9%"],
            ["Rice", "0.0978", "79.1%"],
            ["Tomatoes", "0.0432", "82.5%"],
            ["All", "per-commodity", "84.0%"],
        ],
        0.75,
        3.78,
        5.75,
        1.85,
        font_size=8.6,
    )
    d.image(s, ROOT / "visualizations" / "fig_calibration_coverage.png", 6.85, 3.7, 5.45, 2.1)
    d.text(s, "Plain meaning: CQR widens the predicted band only as much as the 2022 validation errors say is needed.", 1.0, 6.15, 11.2, 0.34, size=12.0, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)

    # 10
    s = d.slide("Evaluation Protocol", "Evaluation")
    d.table(
        s,
        ["Metric", "Formula / definition", "Interpretation"],
        [
            ["MAE", "mean |y - q50|", "Average Rs/KG error"],
            ["MAPE", "mean |(y - q50)/y| x 100", "Scale-free percent error"],
            ["Coverage", "mean 1[q10 <= y <= q90]", "Band reliability"],
            ["Band width", "mean(q90 - q10)", "Usefulness of uncertainty"],
        ],
        0.75,
        1.08,
        6.15,
        2.1,
        font_size=8.8,
    )
    d.code(
        s,
        "Test set discipline:\n"
        "Train final TFT     : through 2021-12\n"
        "Calibration window  : 2022\n"
        "Held-out test       : 2023+\n"
        "Matched test rows   : n = 792\n\n"
        "Statistical checks:\n"
        "paired t-test, Wilcoxon, Cohen's d,\n"
        "binomial coverage test, VSN bootstrap.",
        7.25,
        1.08,
        5.1,
        2.28,
        size=8.8,
    )
    d.stat(s, "Ablation rows", "792", 0.78, 4.02, accent="blue")
    d.stat(s, "Bootstrap resamples", "1000", 3.22, 4.02, accent="violet")
    d.stat(s, "Coverage target", "80%", 5.66, 4.02, accent="green")
    d.stat(s, "Horizon", "1-6 mo", 8.1, 4.02, accent="amber")
    d.stat(s, "Encoder", "24 mo", 10.54, 4.02, accent="teal")
    d.bullets(
        s,
        [
            "XGBoost is reported honestly as the better point-only baseline.",
            "TFT-XGBFusion-CQR is evaluated as the best combined system: point + band + explanation.",
            "Coverage is empirical test coverage, not a theoretical claim copied from quantile labels.",
        ],
        0.9,
        5.35,
        11.2,
        1.05,
        size=11.2,
    )

    # 11
    s = d.slide("Main Results on 2023+ Test Set", "Results")
    d.table(
        s,
        ["Model", "MAE", "MAPE", "Coverage", "Band width"],
        [
            ["XGBoost point baseline", "1.95", "3.5%", "--", "--"],
            ["TFT-Base", "10.53", "29.1%", "64.9%", "23.62"],
            ["TFT-EnsCQR", "10.85", "28.3%", "88.5%", "32.44"],
            ["TFT-Retrain21-CQR", "9.91", "29.2%", "87.6%", "34.78"],
            ["TFT-XGBFusion-CQR", "5.27", "11.4%", "84.0%", "11.99"],
        ],
        0.72,
        1.1,
        7.0,
        2.35,
        font_size=8.9,
    )
    d.stat(s, "Final MAE", "5.27", 8.15, 1.12, accent="teal")
    d.stat(s, "Final MAPE", "11.4%", 10.55, 1.12, accent="green")
    d.stat(s, "Coverage", "84.0%", 8.15, 2.25, accent="blue")
    d.stat(s, "Band width", "11.99", 10.55, 2.25, accent="violet")
    d.image(s, ROOT / "visualizations" / "fig_ablation_mae.png", 0.72, 3.95, 5.65, 2.35)
    d.image(s, ROOT / "Report&Paper" / "figures" / "evaluation_metrics.png", 6.72, 3.95, 5.65, 2.35)
    d.text(s, "Result story: XGBoost wins raw point error; the final TFT stack wins the complete forecasting objective with calibrated uncertainty and inspectable explanations.", 0.95, 6.55, 11.3, 0.28, size=11.0, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)

    # 12
    s = d.slide("Commodity-Level Behavior", "Results")
    d.table(
        s,
        ["Commodity", "MAE", "MAPE", "Coverage", "Interpretation"],
        [
            ["Onions", "2.25", "9.1%", "91%", "Best calibrated; seasonal memory helps"],
            ["Tomatoes", "8.13", "13.0%", "82%", "Hardest due to sudden supply shocks"],
            ["Rice", "4.81", "11.7%", "79%", "Stable series, narrower uncertainty need"],
        ],
        0.75,
        1.05,
        6.4,
        1.65,
        font_size=8.6,
    )
    d.image(s, ROOT / "visualizations" / "quantile_forecast_tomatoes.png", 0.75, 3.08, 3.85, 2.7)
    d.image(s, ROOT / "visualizations" / "quantile_forecast_onions.png", 4.82, 3.08, 3.85, 2.7)
    d.image(s, ROOT / "visualizations" / "quantile_forecast_rice.png", 8.89, 3.08, 3.45, 2.7)
    d.bullets(
        s,
        [
            "Tomatoes show the widest practical risk because monthly data cannot capture short one-week crop failures.",
            "Onions benefit from annual lag and rolling features around known crisis seasons.",
            "Rice checks whether the same architecture can avoid exaggerating stable commodity risk.",
        ],
        7.55,
        1.04,
        4.72,
        1.38,
        size=10.0,
    )

    # 13
    s = d.slide("Explainability: Stable Feature Rankings", "Explainability")
    d.stat(s, "Kendall tau", "0.942", 0.75, 1.05, accent="green")
    d.stat(s, "95% CI low", "0.895", 3.15, 1.05, accent="blue")
    d.stat(s, "95% CI high", "0.990", 5.55, 1.05, accent="violet")
    d.image(s, ROOT / "visualizations" / "fig_vsn_stability.png", 0.75, 2.28, 5.7, 3.3)
    d.image(s, ROOT / "visualizations" / "variable_importance_comparison.png", 6.75, 2.28, 5.5, 3.3)
    d.table(
        s,
        ["Top encoder feature", "Mean VSN weight"],
        [
            ["temperature_mean", "0.143"],
            ["rain_excess", "0.143"],
            ["rain_deficit", "0.140"],
            ["humidity_mean", "0.066"],
            ["year", "0.060"],
        ],
        8.25,
        1.05,
        4.0,
        1.02,
        font_size=8.2,
    )
    d.text(s, "The explanation is not just visual: ranking stability is tested with 1000 bootstrap resamples across series.", 1.0, 6.1, 11.0, 0.33, size=11.5, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)

    # 14
    s = d.slide("Attention and Forecast Evidence", "Explainability")
    d.image(s, ROOT / "visualizations" / "attention_heatmap.png", 0.75, 1.08, 5.75, 2.55)
    d.image(s, ROOT / "visualizations" / "attention_distribution.png", 6.85, 1.08, 5.45, 2.55)
    d.code(
        s,
        "Interpretation rule:\n"
        "High attention at t-1  -> recent momentum\n"
        "High attention at t-12 -> annual seasonality\n"
        "Spread attention       -> uncertain regime\n\n"
        "VSN explains variables; attention explains time.",
        0.78,
        4.18,
        5.72,
        1.55,
        size=9.0,
    )
    d.bullets(
        s,
        [
            "Dynamic VSN weights can change per market, commodity, and month.",
            "Attention weights expose which encoder months affect each decoder horizon.",
            "Together they let the dashboard answer why a forecast moved, not just what number changed.",
        ],
        6.95,
        4.2,
        5.2,
        1.25,
        size=11.0,
    )

    # 15
    s = d.slide("Dashboard: Four Operational Views", "Implementation")
    d.image(s, ROOT / "Report&Paper" / "figures" / "dashboard_view1_forecast.png", 0.75, 1.05, 5.7, 2.25)
    d.image(s, ROOT / "Report&Paper" / "figures" / "dashboard_view2_future.png", 6.75, 1.05, 5.7, 2.25)
    d.image(s, ROOT / "Report&Paper" / "figures" / "dashboard_view3_explainability.png", 0.75, 3.8, 5.7, 2.25)
    d.image(s, ROOT / "Report&Paper" / "figures" / "dashboard_view4_evaluation.png", 6.75, 3.8, 5.7, 2.25)
    d.text(s, "Views: historical forecast, future forecast, model explainability, and present-vs-predicted test validation.", 1.0, 6.42, 11.0, 0.32, size=11.2, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)

    # 16
    s = d.slide("Reproducible Engineering Map", "Implementation")
    d.table(
        s,
        ["Stage", "Script / file", "Main output"],
        [
            ["Filter prices", "00_filter_prices.py", "prices_filtered.csv"],
            ["Weather", "01_fetch_weather.py", "nasa_weather_1994_2026.csv"],
            ["Merge", "02_merge_features.py", "master_dataset.csv"],
            ["TFT", "03_train_tft.py", "tft_best.ckpt"],
            ["XGBoost", "04_train_xgboost.py", "xgb_baseline.pkl"],
            ["Fusion", "10_xgb_as_tft_feature.py", "master_dataset_xgbfused.csv"],
            ["Calibrate", "08_conformal_calibrate.py", "conformal_offsets_step5.json"],
            ["Validate", "11_explainability_stats.py", "statistics + figures"],
        ],
        0.75,
        1.04,
        6.9,
        3.5,
        font_size=8.2,
    )
    d.code(
        s,
        "Run dashboard:\nstreamlit run app.py\n\n"
        "Rebuild evaluation:\npython scripts/06_evaluate.py\npython scripts/11_explainability_stats.py\n\n"
        "Best family auto-selection:\nstep5 -> step1 -> original",
        8.05,
        1.08,
        4.15,
        2.0,
        size=9.1,
    )
    d.bullets(
        s,
        [
            "No separate model server is required for the class demo.",
            "Artifacts are loaded from models/ and data/processed/.",
            "Dashboard fallback keeps the app usable even if the fused checkpoint is missing.",
            "Generated figures are reused consistently in paper, report, README, and PPT.",
        ],
        8.1,
        3.65,
        4.05,
        1.55,
        size=10.4,
    )

    # 17
    s = d.slide("Limitations and Future Work", "Closing")
    d.table(
        s,
        ["Current limitation", "Effect", "Future improvement"],
        [
            ["Monthly data", "Weekly spikes are smoothed", "Daily/weekly AGMARKNET expansion"],
            ["Monthly weather means", "Short heat shocks invisible", "Higher frequency weather features"],
            ["Three commodities", "Limited generality test", "Extend to pulses, oils, vegetables"],
            ["Static calibration", "Regime shifts may drift", "Rolling conformal recalibration"],
            ["Prototype dashboard", "Manual artifact refresh", "Scheduled data/model update job"],
        ],
        0.76,
        1.1,
        7.0,
        2.65,
        font_size=8.4,
    )
    d.card(s, 8.25, 1.14, 3.85, 2.62)
    d.text(s, "Final contribution", 8.55, 1.44, 3.2, 0.25, size=13.0, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)
    d.bullets(
        s,
        [
            "Leakage-controlled XGBoost-to-TFT fusion",
            "Probabilistic q10/q50/q90 forecasting",
            "Per-commodity CQR calibration",
            "Statistically tested explanations",
            "Usable Streamlit decision interface",
        ],
        8.55,
        1.92,
        3.1,
        1.35,
        size=9.5,
    )
    d.text(s, "Conclusion: the system turns volatile food-price prediction into a decision-support workflow: expected price, realistic band, and model evidence on the same screen.", 1.0, 4.65, 11.1, 0.55, size=15.0, bold=True, color=COLORS["ink"], align=PP_ALIGN.CENTER)
    d.text(s, "Thank you", 4.4, 5.75, 4.4, 0.45, size=24, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)

    # 18
    s = d.slide("Key References", "Closing")
    d.bullets(
        s,
        [
            "Lim et al. (2021): Temporal Fusion Transformers for interpretable multi-horizon time series forecasting.",
            "Romano et al. (2019): Conformalized Quantile Regression for distribution-free prediction intervals.",
            "World Food Programme India food price dataset, HDX.",
            "NASA POWER monthly point API for rainfall, temperature, and humidity covariates.",
            "Project report and IEEE-style paper in Report&Paper/ are the source documents for this deck.",
        ],
        0.85,
        1.18,
        7.2,
        2.6,
        size=12.0,
    )
    d.code(
        s,
        "Core claim:\n"
        "TFT-XGBFusion-CQR = point signal + quantile model + conformal reliability + built-in explanation.\n\n"
        "This is why the final system is not only a model comparison; it is a complete forecasting workflow.",
        8.35,
        1.25,
        3.8,
        1.85,
        size=9.5,
    )
    d.image(s, ROOT / "visualizations" / "feature_importance_xgboost.png", 0.85, 4.22, 5.6, 2.15)
    d.image(s, ROOT / "visualizations" / "tft_encoder_importance.png", 6.85, 4.22, 5.3, 2.15)

    d.prs.save(OUTPUT_PPTX)
    print(f"Saved {OUTPUT_PPTX}")


if __name__ == "__main__":
    build_deck()
